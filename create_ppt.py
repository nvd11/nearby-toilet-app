from pptx import Presentation

prs = Presentation()

# Slide 1: Section 1
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "第一部分：市场分析"
slide.placeholders[1].text = "2026年Q1市场回顾"

slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "市场分析详情"
slide.placeholders[1].text = "这里是第一部分的详细内容"

# Slide 2: Section 2
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "第二部分：产品规划"
slide.placeholders[1].text = "Q2产品迭代路线图"

slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "规划详情"
slide.placeholders[1].text = "这里是第二部分的详细内容"

prs.save('/home/gateman/.openclaw/workspace/Section_Demo.pptx')
