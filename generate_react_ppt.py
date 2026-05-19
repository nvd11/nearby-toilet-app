from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

# Slide 1: Title
slide_layout = prs.slide_layouts[0] 
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "ReAct Agent 核心原理解析"
subtitle.text = "让大模型学会思考与行动\n\n分享人：Jason"

# Slide 2: What is ReAct?
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "1. 什么是 ReAct？"
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
tf.text = "ReAct = Reasoning (推理) + Acting (行动)"
p = tf.add_paragraph()
p.text = "传统的 AI 只能被动地根据内置知识回答问题（纯思考）。"
p.level = 1
p = tf.add_paragraph()
p.text = "ReAct 模式赋予了 AI “思考 + 动手” 的能力。"
p.level = 1
p = tf.add_paragraph()
p.text = "核心理念：让 AI 像人类一样，先分析问题，再决定使用什么工具，最后根据工具反馈继续思考。"
p.level = 1

# Slide 3: The Core Loop
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "2. ReAct 的核心运转机制 (The Loop)"
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
tf.text = "经典的“三步循环”架构："
p = tf.add_paragraph()
p.text = "🤔 Thought (思考): AI 分析当前面临的问题，决定下一步需要做什么。"
p.level = 1
p = tf.add_paragraph()
p.text = "🛠️ Action (行动): AI 决定调用哪个外部工具（如：网络搜索、查数据库、执行代码）。"
p.level = 1
p = tf.add_paragraph()
p.text = "👀 Observation (观察): 外部工具返回结果，AI 观察这个结果，决定是继续循环还是输出最终答案。"
p.level = 1

# Slide 4: Why it matters?
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "3. 为什么我们需要 ReAct？"
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
tf.text = "打破知识孤岛"
p = tf.add_paragraph()
p.text = "解决 LLM 数据过时导致的“幻觉”，能够实时检索最新信息。"
p.level = 1
p = tf.add_paragraph()
p.text = "复杂任务拆解与自纠错"
p.level = 0
p = tf.add_paragraph()
p.text = "面对需要多步解决的问题，可以一步步执行，遇到错误能自我纠正。"
p.level = 1
p = tf.add_paragraph()
p.text = "连接真实世界"
p.level = 0
p = tf.add_paragraph()
p.text = "让大模型从“聊天机器”进化为可操作 API、控制系统的“智能助理”。"
p.level = 1

# Slide 5: Example
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "4. 一个生动的例子：买电脑"
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
tf.text = "用户提问：帮我对比昨晚发布的苹果新MacBook和上一代的参数。"
p = tf.add_paragraph()
p.text = "Thought 1: 我需要知道昨晚发布了什么。"
p.level = 1
p = tf.add_paragraph()
p.text = "Action 1: [Web Search] 最新苹果发布会MacBook参数。"
p.level = 1
p = tf.add_paragraph()
p.text = "Observation 1: 获得了 M4 芯片 MacBook 的参数。"
p.level = 1
p = tf.add_paragraph()
p.text = "Thought 2: 我还需要上一代 M3 的参数对比。"
p.level = 1
p = tf.add_paragraph()
p.text = "Action 2: [Search] M3 MacBook 参数。"
p.level = 1
p = tf.add_paragraph()
p.text = "Thought 3/Action 3: 信息充足，生成最终对比建议。"
p.level = 1

prs.save("/home/gateman/.openclaw/workspace/ReAct_Agent_Sharing.pptx")
print("PPT Generated Successfully!")
